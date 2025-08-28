from ft_scraper.load.db import get_db_connection,get_recent_articles
from tqdm import tqdm
import time
from google import genai
import os
from dotenv import load_dotenv
import json
from collections import defaultdict
from sklearn.cluster import KMeans
from sklearn.metrics import silhouette_score
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt
from wordcloud import WordCloud
import matplotlib.pyplot as plt
import io

from datetime import datetime

def get_embeddings(texts,API_KEY):
    try:
        
        if not API_KEY:
            raise ValueError("GEMINI_API_KEY not found in .env file")
        
        client = genai.Client(api_key=API_KEY)
        result = client.models.embed_content(
            model="gemini-embedding-001",
            contents=texts
        )
        return [e.values for e in result.embeddings]
    
    except Exception as ex:
        print("Error:", ex)
        return None

def batch_embeddings(texts,API_KEY, batch_size=5):
    all_embeddings = []
    
    # Wrap the range with tqdm for progress bar
    for i in tqdm(range(0, len(texts), batch_size), desc="Embedding batches"):
        batch = texts[i:i+batch_size]
        embeddings = get_embeddings(batch,API_KEY)
        time.sleep(10)  # optional delay between batches
        all_embeddings.extend(embeddings)
    
    return all_embeddings

def choose_optimal_k(embeddings, k_min=2, k_max=10):
    best_k = k_min
    best_score = -1

    for k in range(k_min, min(k_max, len(embeddings)) + 1):
        kmeans = KMeans(n_clusters=k, random_state=42, n_init=10)
        labels = kmeans.fit_predict(embeddings)
        score = silhouette_score(embeddings, labels)

        if score > best_score:
            best_k = k
            best_score = score

    return best_k

def summarize_theme(articles, model_client, max_tokens=300):
    """
    Summarize a list of articles belonging to the same theme into one collective theme summary.
    Always return structured JSON.
    """
    combined_text = "\n\n".join(articles)
    
    prompt = f"""
        You will read multiple news articles on the same theme. 
        Produce ONE integrated summary in **strict JSON** format with the following keys:

        {{
            "headline": "string, one sentence, concise like a newspaper headline",
            "main_idea": "string, 2-3 sentences capturing the overarching theme",
            "subtopics": [
                "string, distinct aspect 1",
                "string, distinct aspect 2",
                "string, distinct aspect 3"
            ]
        }}

        Do NOT summarize each article separately. Only ONE integrated summary.
        
        Articles:
        {combined_text}
    """
    
    response = model_client.models.generate_content(
        model="gemini-2.0-flash",
        contents=prompt
    )
    
    # Ensure clean JSON response
    import json
    try:
        result = json.loads(response.text.strip())
    except json.JSONDecodeError:
        # fallback: try to extract JSON substring
        import re
        match = re.search(r"\{.*\}", response.text, re.S)
        if match:
            result = json.loads(match.group(0))
        else:
            raise ValueError("Model did not return valid JSON.")
    
    return result

def generate_presentation(presentation_data, output_file=f'data/presentations/themes_presentation_{datetime.now().strftime("%Y-%m-%d_%H-%M")}.pptx'):
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    print("Generate Custom Presentation ...")

    for theme in presentation_data:
        slide_layout = prs.slide_layouts[5]  # Title only
        slide = prs.slides.add_slide(slide_layout)

        # --- Format summary JSON ---
        
        # Add summary box
        left, top, width, height = Inches(0.5), Inches(0.5), Inches(15), Inches(3.5)
        summary_box = slide.shapes.add_textbox(left, top, width, height)
        tf = summary_box.text_frame
        tf.word_wrap = True
        tf.clear()  # start fresh

        summary_data = theme["summary"]

        # Headline
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "Headline: "
        run.font.bold = True
        run.font.size = Pt(14)
        run = p.add_run()
        run.text = summary_data["headline"]
        run.font.size = Pt(14)
        p.space_after = Pt(12)   # <-- add spacing after headline

        # Main Idea
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "Main Idea: "
        run.font.bold = True
        run.font.size = Pt(14)
        run = p.add_run()
        run.text = summary_data["main_idea"]
        run.font.size = Pt(14)
        p.space_after = Pt(12)   # <-- spacing after main idea

        # Subtopics
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "Subtopics:"
        run.font.bold = True
        run.font.size = Pt(14)
        p.space_after = Pt(6)    # <-- small spacing before bullets

        for sub in summary_data.get("subtopics", []):
            p = tf.add_paragraph()
            p.text = f"- {sub}"
            p.level = 1
            p.font.size = Pt(14)
            p.space_after = Pt(6)  # space between bullet points

        # Word cloud
        articles_text = " ".join([a.get("content", a["headline"]) for a in theme["articles"]])
        wc = WordCloud(width=800, height=400, background_color="white").generate(articles_text)

        img_buf = io.BytesIO()
        plt.figure(figsize=(8, 4))
        plt.imshow(wc, interpolation="bilinear")
        plt.axis("off")
        plt.tight_layout()
        plt.savefig(img_buf, format="png", bbox_inches='tight')
        plt.close()
        img_buf.seek(0)

        left, top = Inches(0.50), Inches(3.60)
        slide.shapes.add_picture(img_buf, left, top, width=Inches(5), height=Inches(3.32))

        # References
        left, top, width, height = Inches(0.5), Inches(7.20), Inches(15), Inches(2.0)
        ref_box = slide.shapes.add_textbox(left, top, width, height)
        tf = ref_box.text_frame
        tf.word_wrap = True
        tf.clear()

        # "References:" bold
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "References:"
        run.font.bold = True
        run.font.size = Pt(14)

        for art in theme["articles"][:3]:
            p = tf.add_paragraph()
            
            # Headline normal
            run = p.add_run()
            run.text = f"- {art['headline']} : "
            run.font.size = Pt(14)
            
            # Article ID italic
            run = p.add_run()
            run.text = art["article_id"]
            run.font.italic = True
            run.font.size = Pt(14)

            p.level = 1

    prs.save(output_file)
    print(f"Presentation saved to {output_file}")
    return output_file

def presentation_pipeline():
    
    try:

        load_dotenv(dotenv_path=".env")

        API_KEY = os.getenv("GEMINI_API_KEY")

        # Step 1: Load Articles from MongoDB

        collection = get_db_connection()

        recent_articles  = get_recent_articles(collection)

        # Step 2: Prepare Text for Clustering

        article_texts = []
        article_ids = []

        for article in recent_articles:
            text = " ".join(article.get("content", []))
            article_texts.append(text)
            article_ids.append(article["article_id"])
        
        # Step 3: Convert Text to Embeddings

        embeddings = batch_embeddings(article_texts,API_KEY,batch_size=5)

        # Step 4: Cluster Articles into Themes

        num_themes = choose_optimal_k(embeddings=embeddings)
        kmeans = KMeans(n_clusters=num_themes, random_state=42)
        labels = kmeans.fit_predict(embeddings)

        # Step 5: Aggregate Articles per Theme
        
        themes = defaultdict(list)
        
        for i, label in enumerate(labels):
            themes[label].append({
                "article_id": article_ids[i],
                "headline": recent_articles[i]["topper__headline"],
                "content": article_texts[i]
            })

        # Step 6: Summarize Each Theme
        
        llm_client = genai.Client(api_key=API_KEY)

        summaries = {}  # dict to store all theme summaries

        for theme, articles in themes.items():
            # Extract only the content texts
            article_texts = [a["content"] for a in articles]
            
            theme_summary = summarize_theme(article_texts, llm_client)
            summaries[theme] = theme_summary
        
        # Step 7: Prepare Presentation Data
        
        presentation_data = []

        for theme_id, summary in summaries.items():
            presentation_data.append({
                "theme_id": theme_id,
                "summary": summary,
                "articles": themes[theme_id]  # include headlines or ids for reference
            })

        # Step 8: Generate Presentation

        path = generate_presentation(presentation_data=presentation_data)

        return True,path
    
    except Exception as ex:

        return False,ex